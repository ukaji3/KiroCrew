"""File readers for knowledge ingestion. Supports text, PDF, PPTX, DOCX, HTML."""

import os
import re
from pathlib import Path

from kiro_crew.security import is_sensitive_path

try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

try:
    from docx import Document  # type: ignore[import-untyped]
except ImportError:
    Document = None  # type: ignore[assignment,misc]

try:
    import html2text as _html2text_mod
except ImportError:
    _html2text_mod = None  # type: ignore[assignment]


class FileReader:
    # Binary formats need optional runtime deps: .pdf -> pdfplumber and .docx ->
    # python-docx (both declared in setup.cfg). .pptx -> python-pptx is NOT declared,
    # so .pptx is intentionally kept out of SUPPORTED even though _read_pptx exists.
    SUPPORTED = {
        '', '.md', '.txt', '.org', '.py', '.java', '.ts', '.js', '.rs', '.go',
        '.html', '.htm', '.docx', '.pdf',
        '.csv', '.log', '.json', '.jsonl', '.ndjson', '.yaml', '.yml',
        '.sh', '.rb', '.c', '.cpp', '.h',
    }

    _DISPATCH = {
        '.pdf': '_read_pdf',
        '.pptx': '_read_pptx',
        '.docx': '_read_docx',
        '.html': '_read_html',
        '.htm': '_read_html',
    }

    def read(self, path: str) -> tuple[str, dict]:
        if is_sensitive_path(path):
            raise PermissionError(f"Refusing to read sensitive path: {path}")
        p = Path(path)
        ext = p.suffix.lower()
        base_meta = {
            'format': ext.lstrip('.'),
            'title': p.stem,
            'file_size': os.path.getsize(path),
            'extension': ext,
        }
        method_name = self._DISPATCH.get(ext)
        if method_name:
            text, meta = getattr(self, method_name)(path)
            base_meta.update(meta)
        else:
            text, meta = self._read_text(path, ext.lstrip('.'))
            base_meta.update(meta)
        base_meta['line_count'] = text.count('\n') + 1 if text else 0
        return text, base_meta

    def _read_text(self, path: str, fmt: str) -> tuple[str, dict]:
        try:
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(path, 'r', encoding='latin-1') as f:
                    text = f.read()
            return text, {'format': fmt}
        except Exception as e:
            return f'Error reading file: {e}', {'format': 'error', 'error': str(e)}

    def _read_pdf(self, path: str) -> tuple[str, dict]:
        if pdfplumber is None:
            return ('PDF support requires pdfplumber: pip install pdfplumber',
                    {'format': 'error', 'error': 'PDF support requires pdfplumber'})
        try:
            with pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or '' for p in pdf.pages]
                return '\n'.join(pages), {'format': 'pdf', 'page_count': len(pages)}
        except Exception as e:
            return f'Error reading file: {e}', {'format': 'error', 'error': str(e)}

    def _read_pptx(self, path: str) -> tuple[str, dict]:
        if Presentation is None:
            return ('PPTX support requires python-pptx: pip install python-pptx',
                    {'format': 'error', 'error': 'PPTX support requires python-pptx'})
        try:
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                title = ''
                body_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if shape == slide.shapes.title:
                            title = text
                        else:
                            body_parts.append(text)
                notes = ''
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                section = f'## Slide {i}: {title}\n{chr(10).join(body_parts)}'
                if notes:
                    section += f'\n{notes}'
                parts.append(section)
            return '\n\n'.join(parts), {'format': 'pptx', 'slide_count': len(prs.slides)}
        except Exception as e:
            return f'Error reading file: {e}', {'format': 'error', 'error': str(e)}

    def _read_docx(self, path: str) -> tuple[str, dict]:
        if Document is None:
            return ('DOCX support requires python-docx: pip install python-docx',
                    {'format': 'error', 'error': 'DOCX support requires python-docx'})
        try:
            doc = Document(path)
            lines = []
            for para in doc.paragraphs:
                style = para.style.name if para.style else ''
                text = para.text
                if style.startswith('Heading'):
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    lines.append(f'{"#" * level} {text}')
                else:
                    lines.append(text)
            return '\n'.join(lines), {'format': 'docx', 'content_type': 'markdown', 'paragraph_count': len(doc.paragraphs)}
        except Exception as e:
            return f'Error reading file: {e}', {'format': 'error', 'error': str(e)}

    def _read_html(self, path: str) -> tuple[str, dict]:
        try:
            with open(path, 'r', encoding='utf-8') as f:
                html = f.read()
        except UnicodeDecodeError:
            with open(path, 'r', encoding='latin-1') as f:
                html = f.read()
        except Exception as e:
            return f'Error reading file: {e}', {'format': 'error', 'error': str(e)}
        if _html2text_mod is not None:
            h = _html2text_mod.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            return h.handle(html), {'format': 'html'}
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text, {'format': 'html'}
